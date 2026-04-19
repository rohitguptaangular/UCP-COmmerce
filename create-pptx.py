#!/usr/bin/env python3
"""Generate two PowerPoint decks from the leadership HTML presentations."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──
BG_DARK = RGBColor(0x0F, 0x12, 0x25)
BG_CARD = RGBColor(0x1A, 0x1F, 0x36)
WIPRO = RGBColor(0x5A, 0x2F, 0xC2)
WIPRO_LT = RGBColor(0x7C, 0x4D, 0xFF)
ACCENT = RGBColor(0xE3, 0x00, 0x0F)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
BLUE = RGBColor(0x42, 0xA5, 0xF5)
ORANGE = RGBColor(0xFF, 0x98, 0x00)
WHITE = RGBColor(0xE8, 0xEA, 0xF6)
GRAY = RGBColor(0x9F, 0xA8, 0xDA)
MUTED = RGBColor(0x6B, 0x74, 0xA8)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))


def add_text(tf, text, size=14, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, space_after=4, font_name="Calibri"):
    p = tf.add_paragraph() if len(tf.paragraphs) > 0 and tf.paragraphs[0].text != "" else tf.paragraphs[0]
    if len(tf.paragraphs) > 0 and tf.paragraphs[0].text != "" and p == tf.paragraphs[0]:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(space_after)
    return p


def add_bullet(tf, text, size=12, color=GRAY, bullet_color=None, space_after=3, level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.level = level
    p.space_after = Pt(space_after)
    return p


def add_shape_box(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def slide_header(slide, num, title, subtitle=""):
    set_slide_bg(slide, BG_DARK)
    # Slide number
    tb = add_textbox(slide, 0.7, 0.4, 5, 0.3)
    add_text(tb.text_frame, f"SLIDE {num:02d}", size=9, bold=True, color=WIPRO_LT, space_after=0)
    # Title
    tb = add_textbox(slide, 0.7, 0.7, 8.5, 0.9)
    add_text(tb.text_frame, title, size=28, bold=True, color=WHITE, space_after=4)
    # Subtitle
    if subtitle:
        tb = add_textbox(slide, 0.7, 1.55, 8, 0.6)
        add_text(tb.text_frame, subtitle, size=13, color=MUTED, space_after=0)


def add_table(slide, left, top, width, rows_data, col_widths=None):
    """rows_data: list of lists. First row = header."""
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(0.35 * rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for r, row in enumerate(rows_data):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.text = cell_text
            p.font.size = Pt(10 if r == 0 else 11)
            p.font.bold = (r == 0)
            p.font.color.rgb = MUTED if r == 0 else GRAY
            p.font.name = "Calibri"

            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_CARD if r == 0 else BG_DARK

    return table_shape


# ============================================================
#  DECK 1: Technical Leadership Deck (from leadership-deck.html)
# ============================================================

def create_deck1():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── SLIDE 1: Title ──
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG_DARK)

    # Wipro logo box
    shape = add_shape_box(slide, 0.7, 0.6, 0.55, 0.55, fill_color=WIPRO)
    shape.text_frame.paragraphs[0].text = "W"
    shape.text_frame.paragraphs[0].font.size = Pt(20)
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.color.rgb = WHITE
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.word_wrap = True

    tb = add_textbox(slide, 0.7, 2.2, 10, 1.5)
    add_text(tb.text_frame, "Wipro Agentic Commerce Accelerator", size=36, bold=True, color=WHITE, space_after=8)
    add_text(tb.text_frame, "Connecting Enterprise Commerce Platforms to AI Shopping Agents", size=18, color=GRAY, space_after=16)
    add_text(tb.text_frame, "Technical Architecture & Solution Overview  |  Confidential", size=12, color=MUTED)

    # ── SLIDE 2: The Opportunity ──
    slide = prs.slides.add_slide(blank)
    slide_header(slide, 1, "The Opportunity: AI Is the New Storefront",
                 "Google, OpenAI, and Microsoft launched AI shopping protocols in Jan-Feb 2026. Merchants who aren't connected are invisible.")

    # Three cards
    cards = [
        ("Google UCP", "Jan 2026", "Universal Commerce Protocol. Buy directly inside Gemini & Google Search AI Mode. Payment via Google Pay.", ACCENT),
        ("OpenAI ACP", "Feb 2026", '"Buy it in ChatGPT" with Stripe. 500M+ ChatGPT users can purchase without leaving the conversation.', BLUE),
        ("Microsoft Copilot", "Jan 2026", "Copilot Checkout with Shopify. Shopping through Copilot integrated into Windows, Edge, Office.", WIPRO_LT),
    ]
    for i, (title, date, desc, clr) in enumerate(cards):
        x = 0.7 + i * 4.1
        box = add_shape_box(slide, x, 2.3, 3.8, 2.0, border_color=clr)
        tf = box.text_frame
        tf.word_wrap = True
        add_text(tf, title, size=15, bold=True, color=WHITE, space_after=2)
        add_text(tf, date, size=10, bold=True, color=clr, space_after=6)
        add_text(tf, desc, size=11, color=GRAY, space_after=0)

    # Threat callout
    box = add_shape_box(slide, 0.7, 4.6, 11.9, 0.8, border_color=ACCENT)
    tf = box.text_frame
    tf.word_wrap = True
    add_text(tf, 'When a customer asks AI to buy a product, AI only shows merchants who are connected. If our clients aren\'t connected, AI recommends competitors instead.', size=13, bold=True, color=WHITE, space_after=0)

    # Era table
    era_data = [
        ["Era", "How Customers Shop", "What Merchants Need"],
        ["2000s", "Type URLs in browser", "A website"],
        ["2010s", "Tap apps on phone", "A mobile app"],
        ["2020s", "Search on marketplaces", "Marketplace presence"],
        ["2026+", "Talk to AI assistants", "AI-compatible commerce APIs"],
    ]
    add_table(slide, 0.7, 5.6, 11.9, era_data, col_widths=[1.5, 5, 5.4])

    # ── SLIDE 3: What Wipro Builds ──
    slide = prs.slides.add_slide(blank)
    slide_header(slide, 2, "What Wipro Builds: A Thin Translation Layer",
                 "Not middleware. Not a platform. A lightweight API that makes existing e-commerce stores visible to AI agents.")

    # Architecture box
    box = add_shape_box(slide, 0.7, 2.2, 11.9, 4.5, border_color=RGBColor(0x2A, 0x2F, 0x48))
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].space_before = Pt(4)
    lines = [
        ("Google Gemini  /  ChatGPT  /  Microsoft Copilot", BLUE),
        ("        |                    |                    |", MUTED),
        ("        |  (they call us)    |  (they call us)    |  (they call us)", MUTED),
        ("        v                    v                    v", MUTED),
        ("", WHITE),
        ("   WIPRO ACCELERATOR  (hosted on client's cloud)", WIPRO_LT),
        ("", WHITE),
        ("   Receives request  ->  Translates fields  ->  Calls SAP", WIPRO_LT),
        ("   Gets SAP response ->  Translates back    ->  Returns to AI", WIPRO_LT),
        ("", WHITE),
        ("   No business logic. No pricing. No decisions.", WIPRO_LT),
        ("   ~500 lines of code. 5-6 endpoints.", WIPRO_LT),
        ("", WHITE),
        ("        |  (calls existing SAP REST APIs)", MUTED),
        ("        v", MUTED),
        ("", WHITE),
        ("   CLIENT'S SAP COMMERCE CLOUD  (completely untouched)", GREEN),
        ("   Pricing, inventory, promotions, tax, shipping, fulfillment", GREEN),
    ]
    for txt, clr in lines:
        p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(11)
        p.font.color.rgb = clr
        p.font.name = "Courier New"
        p.space_after = Pt(1)

    # ── SLIDE 4: Does vs Doesn't ──
    slide = prs.slides.add_slide(blank)
    slide_header(slide, 3, "What the Accelerator Does vs Does NOT Do",
                 "Zero risk to the client's existing commerce platform. We only translate and forward.")

    # Left card: Does
    box = add_shape_box(slide, 0.7, 2.2, 5.7, 4.5, border_color=GREEN)
    tf = box.text_frame
    tf.word_wrap = True
    add_text(tf, "What It DOES", size=16, bold=True, color=GREEN, space_after=8)
    does = [
        "Expose UCP / ACP / MCP API endpoints",
        "Translate field names between AI platforms and SAP",
        "Route requests to SAP OCC REST APIs",
        "Track shopping session-to-cart mapping",
        "Handle customer account linking (OAuth2)",
        "Return responses in each AI platform's expected format",
    ]
    for d in does:
        add_bullet(tf, d, size=12, color=GRAY, space_after=4)

    # Right card: Doesn't
    box = add_shape_box(slide, 6.9, 2.2, 5.7, 4.5, border_color=ACCENT)
    tf = box.text_frame
    tf.word_wrap = True
    add_text(tf, "What It Does NOT Do", size=16, bold=True, color=ACCENT, space_after=8)
    doesnt = [
        "Calculate pricing (SAP does this)",
        "Apply discounts or promotions (SAP does this)",
        "Validate inventory or stock (SAP does this)",
        "Determine shipping costs (SAP does this)",
        "Store customer or payment data",
        "Make any business decisions whatsoever",
    ]
    for d in doesnt:
        add_bullet(tf, d, size=12, color=GRAY, space_after=4)

    # ── SLIDE 5: Wipro's Role & IP ──
    slide = prs.slides.add_slide(blank)
    slide_header(slide, 4, "Wipro's Role & Reusable IP",
                 "What Wipro owns, delivers, and reuses across every client engagement.")

    roles = [
        ("1", "Build the Accelerator", "Develop the thin API layer using reusable templates. ~500 lines of code."),
        ("2", "Maintain Protocol Compliance", "Update specs as Google/OpenAI change protocols. ~0.5 FTE ongoing."),
        ("3", "Deploy on Client's Cloud", "Same cloud as their SAP (Azure/AWS/GCP). Private networking. ~$1K/mo."),
        ("4", "Register with AI Platforms", "Google Merchant Center, OpenAI portal, MCP config. Form filling."),
        ("5", "Test with AI Simulators", "Automated Gemini/ChatGPT bots test the full shopping flow."),
        ("6", "Operate (99.9% SLA)", "24/7 monitoring, alerting, incident response. Production-grade."),
        ("7", "Analytics Dashboard", "Which AI drives traffic, conversion rates, popular products via AI."),
        ("8", "Scale to New Clients", "Same kit, new config. Each new client = faster and cheaper."),
    ]

    for i, (num, title, desc) in enumerate(roles):
        col = i % 4
        row = i // 4
        x = 0.7 + col * 3.15
        y = 2.2 + row * 2.2
        box = add_shape_box(slide, x, y, 2.95, 1.9, border_color=WIPRO)
        tf = box.text_frame
        tf.word_wrap = True
        add_text(tf, f"{num}. {title}", size=12, bold=True, color=WIPRO_LT, space_after=4)
        add_text(tf, desc, size=10, color=GRAY, space_after=0)

    # ── SLIDE 6: Reusable IP Kit ──
    slide = prs.slides.add_slide(blank)
    slide_header(slide, 5, "The Accelerator Kit (Wipro's Reusable IP)",
                 "Built once, reused for every client. The knowledge compounds with each engagement.")

    kit_items = [
        ("API Specifications", "OpenAPI definitions for UCP, ACP, MCP. Maintained as protocols evolve. Client doesn't read 200-page specs."),
        ("Canonical Data Model", "Standardized internal field names between all protocols and backends. Add new protocol = +1 mapping, not N."),
        ("Mapping Sheets", "Field-by-field transformation rules. The core IP. 'city' vs 'town' vs 'locality' solved once."),
        ("Flow Templates", "Orchestration blueprints: create cart > add items > address > checkout. Platform-agnostic."),
        ("Config Templates", "Discovery profile generator, deployment manifests. New client = fill config, deploy."),
        ("Test Suite", "Protocol compliance tests + AI agent simulator. Validate end-to-end in hours."),
    ]

    for i, (title, desc) in enumerate(kit_items):
        col = i % 3
        row = i // 3
        x = 0.7 + col * 4.15
        y = 2.2 + row * 2.4
        box = add_shape_box(slide, x, y, 3.9, 2.1, border_color=WIPRO)
        tf = box.text_frame
        tf.word_wrap = True
        add_text(tf, f"{i+1}. {title}", size=13, bold=True, color=WIPRO_LT, space_after=6)
        add_text(tf, desc, size=10, color=GRAY, space_after=0)

    # Reusability row
    reuse = [
        ("Client 1 (SAP)", "4 weeks", "Build from scratch"),
        ("Client 2 (SFCC)", "2-3 weeks", "New backend connector"),
        ("Client 3 (SAP)", "1-2 weeks", "Config change only"),
        ("Client 4+ (Any)", "1-2 weeks", "Config change only"),
    ]
    for i, (client, time, note) in enumerate(reuse):
        x = 0.7 + i * 3.15
        box = add_shape_box(slide, x, 5.0, 2.95, 1.2, border_color=GREEN if i >= 2 else BLUE)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = time
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = GREEN if i >= 2 else BLUE
        p.alignment = PP_ALIGN.CENTER
        add_text(tf, client, size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER, space_after=1)
        add_text(tf, note, size=9, color=MUTED, alignment=PP_ALIGN.CENTER, space_after=0)

    # ── SLIDE 7: Deployment ──
    slide = prs.slides.add_slide(blank)
    slide_header(slide, 6, "Deployment: Where & How",
                 "Deployed on the client's own cloud, same network as their SAP Commerce.")

    deploy_data = [
        ["If Client's SAP Is On", "We Deploy On", "Cloud Services Used"],
        ["Microsoft Azure", "Azure (same subscription)", "Container Apps + API Mgmt + Redis Cache"],
        ["Amazon AWS", "AWS (same account)", "ECS Fargate + API Gateway + ElastiCache"],
        ["Google Cloud", "GCP (same project)", "Cloud Run + Apigee + Memorystore"],
        ["SAP BTP", "SAP BTP (same subaccount)", "Kyma Runtime + SAP API Management"],
    ]
    add_table(slide, 0.7, 2.2, 11.9, deploy_data, col_widths=[3.0, 4.0, 4.9])

    # Ownership
    own_data = [
        ["Wipro Delivers & Manages", "Client Provides", "Already Exists (Untouched)"],
        ["Build the accelerator layer", "Cloud subscription access", "SAP Commerce Cloud"],
        ["Deploy on client's cloud", "Domain (agentic.client.com)", "Product catalog & pricing"],
        ["Register with AI platforms", "SAP API credentials", "Order & fulfillment systems"],
        ["Operate with 99.9% SLA", "Google Merchant Center acct", "Payment processing"],
        ["Update when protocols change", "Payment provider accounts", "Customer CRM"],
        ["Analytics dashboard", "Brand/product information", "Existing website & app"],
    ]
    add_table(slide, 0.7, 4.0, 11.9, own_data, col_widths=[4.0, 4.0, 3.9])

    # ── SLIDE 8: Q&A ──
    slide = prs.slides.add_slide(blank)
    slide_header(slide, 7, "Anticipated Questions",
                 "Answers to the questions leadership will ask.")

    qa_data = [
        ["Question", "Answer"],
        ["Do we integrate with each AI platform?", "No. We expose endpoints, they call us. We register on each portal (form filling, not coding)."],
        ["Do we change anything in SAP?", "Zero changes. We only call existing SAP REST APIs. No config, no logic, no data changes."],
        ["Is this secure?", "Payment NEVER touches our layer. Google Pay/Stripe handle it tokenized. No PII stored."],
        ["Can't SAP build this natively?", "They will — in 12-18 months. Clients need it NOW. We bridge the gap."],
        ["Why not MuleSoft?", "If client has it, we use it. Otherwise, $100K/yr license for 5 endpoints is overkill."],
        ["What if protocols change?", "That's the maintenance retainer. We update mappings & endpoints. ~0.5 FTE ongoing."],
        ["What's Wipro's IP?", "API specs, canonical data model, mapping sheets, flow templates, test suites. Reusable."],
        ["What if client uses Shopify?", "Shopify is building native UCP. Our sweet spot: SAP & Salesforce Commerce clients."],
    ]
    add_table(slide, 0.7, 2.2, 11.9, qa_data, col_widths=[4.0, 7.9])

    return prs


# ============================================================
#  DECK 2: Business Leadership Deck (from leadership-deck-v2.html)
# ============================================================

def create_deck2():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── TITLE SLIDE ──
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG_DARK)
    shape = add_shape_box(slide, 0.7, 0.6, 0.55, 0.55, fill_color=WIPRO)
    shape.text_frame.paragraphs[0].text = "W"
    shape.text_frame.paragraphs[0].font.size = Pt(20)
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.color.rgb = WHITE
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    tb = add_textbox(slide, 0.7, 2.2, 10, 1.8)
    add_text(tb.text_frame, "Wipro Agentic Commerce Practice", size=36, bold=True, color=WHITE, space_after=8)
    add_text(tb.text_frame, "Helping Retail & CPG Clients Sell Through AI Shopping Assistants", size=18, color=GRAY, space_after=16)
    add_text(tb.text_frame, "Leadership Briefing  |  Confidential  |  March 2026", size=12, color=MUTED)

    # ── SLIDE 1: THE SHIFT ──
    slide = prs.slides.add_slide(blank)
    slide_header(slide, 1, "Shopping Is Moving Inside AI Conversations",
                 "In Jan-Feb 2026, Google, OpenAI, and Microsoft each launched ways for customers to buy products directly through AI — without visiting a website.")

    # Three cards
    cards = [
        ("Google", "Jan 2026", "Customers can now buy products directly inside Google Search and the Gemini app by talking to AI.", ACCENT),
        ("OpenAI / ChatGPT", "Feb 2026", '500M+ ChatGPT users can now say "buy me X" and purchase without leaving the conversation.', BLUE),
        ("Microsoft Copilot", "Jan 2026", "Shopping through Copilot built into Windows, Edge, and Office. Partnership with Shopify.", WIPRO_LT),
    ]
    for i, (title, date, desc, clr) in enumerate(cards):
        x = 0.7 + i * 4.1
        box = add_shape_box(slide, x, 2.2, 3.85, 1.9, border_color=clr)
        tf = box.text_frame
        tf.word_wrap = True
        add_text(tf, title, size=15, bold=True, color=WHITE, space_after=2)
        add_text(tf, date, size=10, bold=True, color=clr, space_after=6)
        add_text(tf, desc, size=11, color=GRAY, space_after=0)

    # Simple era table
    era_data = [
        ["Era", "How Customers Shop", "What Merchants Needed"],
        ["2000s", "Visit websites", "A website"],
        ["2010s", "Use mobile apps", "A mobile app"],
        ["2020s", "Search on Amazon, Flipkart", "Marketplace presence"],
        ["2026+", 'Tell AI: "Buy me a Dyson V15"', "AI-compatible commerce connection"],
    ]
    add_table(slide, 0.7, 4.4, 11.9, era_data, col_widths=[1.5, 5.2, 5.2])

    # Threat
    box = add_shape_box(slide, 0.7, 6.1, 11.9, 0.7, border_color=ACCENT)
    tf = box.text_frame
    tf.word_wrap = True
    add_text(tf, "THE RISK: If our clients' stores aren't connected to these AI assistants, the AI will recommend and sell competitors' products instead.", size=13, bold=True, color=WHITE, space_after=0)

    # ── SLIDE 2: WHY WIPRO ──
    slide = prs.slides.add_slide(blank)
    slide_header(slide, 2, "Why This Is a Major Opportunity for Wipro",
                 "Every Wipro retail & CPG client needs this. No system integrator has a packaged solution yet. We can be first.")

    # Before/After scenarios
    scenarios = [
        ("WITHOUT Our Solution", ACCENT, [
            ('Customer asks AI: "Buy me Dyson V15"', "Dyson isn't connected. AI says: 'I found similar vacuums from Samsung and Shark.' — LOST SALE."),
            ('Customer asks AI: "Reorder my Nespresso pods"', "Nestle isn't connected. AI redirects to Amazon — LOST MARGIN, LOST CUSTOMER DATA."),
        ]),
        ("WITH Our Solution", GREEN, [
            ('Customer asks AI: "Buy me Dyson V15"', "Dyson IS connected. AI shows V15, checkout in 30 seconds via Google Pay. DIRECT SALE, FULL MARGIN."),
            ('Customer asks AI: "Reorder my Nespresso pods"', "Nestle IS connected. AI knows their account, reorders same pods. RECURRING REVENUE."),
        ]),
    ]

    for i, (heading, clr, items) in enumerate(scenarios):
        x = 0.7 + i * 6.3
        box = add_shape_box(slide, x, 2.2, 5.95, 2.7, border_color=clr)
        tf = box.text_frame
        tf.word_wrap = True
        add_text(tf, heading, size=15, bold=True, color=clr, space_after=8)
        for title, desc in items:
            add_text(tf, title, size=11, bold=True, color=WHITE, space_after=1)
            add_text(tf, desc, size=10, color=GRAY, space_after=8)

    # Why Wipro positioned
    positions = [
        "We manage SAP Commerce for dozens of retail/CPG clients — we know their backends",
        "We manage Salesforce Commerce for another set — same solution works",
        "No SI has a packaged solution yet (Accenture, Infosys, TCS — none)",
        "We already have a working prototype — ready to demo today",
    ]
    box = add_shape_box(slide, 0.7, 5.2, 11.9, 2.0, border_color=WIPRO)
    tf = box.text_frame
    tf.word_wrap = True
    add_text(tf, "Why Wipro Is Perfectly Positioned", size=15, bold=True, color=WIPRO_LT, space_after=8)
    for p in positions:
        add_bullet(tf, p, size=12, color=GRAY, space_after=3)

    # ── SLIDE 3: OUR SOLUTION ──
    slide = prs.slides.add_slide(blank)
    slide_header(slide, 3, "What We Build for the Client",
                 "A thin layer that connects the client's existing online store to AI shopping assistants. We don't change their store.")

    # Restaurant analogy
    box = add_shape_box(slide, 0.7, 2.2, 7.5, 2.2, border_color=GREEN)
    tf = box.text_frame
    tf.word_wrap = True
    add_text(tf, "Simple Analogy: The Restaurant Tablet", size=14, bold=True, color=GREEN, space_after=8)
    add_text(tf, "When a restaurant joins Uber Eats, DoorDash, and Zomato — they don't build 3 kitchens. They add ONE tablet that receives orders from ALL platforms and sends them to the same kitchen.", size=12, color=GRAY, space_after=6)
    add_text(tf, "Our accelerator = the tablet.  Client's SAP = the kitchen.  Google/OpenAI/Microsoft = the delivery apps.", size=12, bold=True, color=WHITE, space_after=0)

    # Emoji visual
    box = add_shape_box(slide, 8.6, 2.2, 4.0, 2.2, border_color=RGBColor(0x2A, 0x2F, 0x48))
    tf = box.text_frame
    tf.word_wrap = True
    add_text(tf, "How It Works", size=14, bold=True, color=WHITE, space_after=10, alignment=PP_ALIGN.CENTER)
    steps = [
        "1. Customer talks to AI assistant",
        "2. AI calls our layer (translates)",
        "3. SAP processes the order",
        "4. Customer gets the product",
    ]
    for s in steps:
        add_text(tf, s, size=11, color=GRAY, space_after=3, alignment=PP_ALIGN.LEFT)

    # What we touch vs don't
    touch_data = [
        ["Wipro Builds (thin layer)", "We Do NOT Touch"],
        ["API endpoints (the 'tablet')", "Client's SAP Commerce configuration"],
        ["Translate AI language to SAP language", "Product catalog or pricing rules"],
        ["Translate SAP responses back to AI", "Discounts, promotions, or tax"],
        ["Track shopping sessions", "Inventory or warehouse systems"],
        ["Handle account linking (OAuth2)", "Payment processing (Google Pay handles this)"],
        ["~500 lines of code. 5-6 endpoints.", "Zero business logic changes. Zero risk."],
    ]
    add_table(slide, 0.7, 4.7, 11.9, touch_data, col_widths=[6.0, 5.9])

    # ── SLIDE 4: WIPRO'S ROLE ──
    slide = prs.slides.add_slide(blank)
    slide_header(slide, 4, "Wipro's Role: End-to-End Ownership",
                 "We don't just build and hand over. We own the full lifecycle — from strategy to production operations.")

    roles = [
        ("1. Assess & Advise", "Evaluate client's platform. Identify which AI platforms to prioritize. Recommend approach.", "Consulting"),
        ("2. Build the Accelerator", "Develop API layer using reusable templates. Configure for client's backend & payment provider.", "Development"),
        ("3. Deploy on Client's Cloud", "Deploy on Azure/AWS/GCP — same cloud as their SAP. Private networking, API gateway, monitoring.", "DevOps"),
        ("4. Register with AI Platforms", "Register on Google Merchant Center, OpenAI portal, publish MCP config. Portal registration, not coding.", "Setup"),
        ("5. Test with Real AI Agents", "Automated Gemini/ChatGPT bots test the full flow: browse, cart, checkout, order tracking.", "QA"),
        ("6. Go Live & Operate", "Production with 99.9% SLA. 24/7 monitoring. Incident response. Client's products now buyable via AI.", "Operations"),
        ("7. Maintain & Evolve", "Update accelerator when Google/OpenAI change protocols. Add new AI platforms as they emerge.", "Ongoing"),
        ("8. Analytics & Insights", "Dashboard: which AI drives most sales, conversion rates, popular products, failed checkouts.", "Value-Add"),
    ]

    for i, (title, desc, tag) in enumerate(roles):
        col = i % 4
        row = i // 4
        x = 0.7 + col * 3.15
        y = 2.2 + row * 2.35
        box = add_shape_box(slide, x, y, 2.95, 2.1, border_color=WIPRO)
        tf = box.text_frame
        tf.word_wrap = True
        add_text(tf, title, size=11, bold=True, color=WIPRO_LT, space_after=4)
        add_text(tf, desc, size=10, color=GRAY, space_after=4)
        add_text(tf, tag, size=9, bold=True, color=MUTED, space_after=0)

    # ── SLIDE 5: DEPLOYMENT ──
    slide = prs.slides.add_slide(blank)
    slide_header(slide, 5, "Deployment & Ownership",
                 "Deployed on the client's own cloud. Client owns the infrastructure. Wipro builds, deploys, and manages it.")

    deploy_data = [
        ["If Client's SAP Is On", "We Deploy On", "Cloud Services We Use"],
        ["Microsoft Azure", "Azure (same subscription)", "Azure Container Apps + API Management + Redis Cache"],
        ["Amazon AWS", "AWS (same account)", "ECS Fargate + API Gateway + ElastiCache"],
        ["Google Cloud", "GCP (same project)", "Cloud Run + Apigee + Memorystore"],
        ["SAP BTP", "SAP BTP (same subaccount)", "Kyma Runtime + SAP API Management"],
    ]
    add_table(slide, 0.7, 2.2, 11.9, deploy_data, col_widths=[3.0, 3.5, 5.4])

    # Note
    tb = add_textbox(slide, 0.7, 4.2, 11.9, 0.4)
    add_text(tb.text_frame, "If client already has MuleSoft, Boomi, or another integration platform — we use that instead. We adapt to the client's existing technology.", size=11, color=MUTED, space_after=0)

    # Ownership columns
    own_headers = ["Wipro Delivers & Manages", "Client Provides", "Already Exists (Untouched)"]
    own_items = [
        ["Build the accelerator layer", "Cloud subscription access", "SAP Commerce Cloud"],
        ["Deploy on client's cloud", "Domain (agentic.client.com)", "Product catalog & pricing"],
        ["Register with AI platforms", "SAP Commerce API credentials", "Order & fulfillment systems"],
        ["Operate with 99.9% SLA", "Google Merchant Center account", "Payment processing (Stripe/Adyen)"],
        ["Update when protocols change", "Payment provider accounts", "Customer CRM"],
        ["Analytics dashboard", "Brand/product information", "Existing website & mobile app"],
        ["Security & compliance", "", ""],
        ["Test with AI agent simulators", "", ""],
    ]
    colors = [WIPRO_LT, GREEN, BLUE]
    for i, (header, clr) in enumerate(zip(own_headers, colors)):
        x = 0.7 + i * 4.15
        box = add_shape_box(slide, x, 4.8, 3.9, 2.6, border_color=clr)
        tf = box.text_frame
        tf.word_wrap = True
        add_text(tf, header, size=12, bold=True, color=clr, space_after=8)
        for row in own_items:
            if row[i]:
                add_bullet(tf, row[i], size=10, color=GRAY, space_after=2)

    # ── SLIDE 6: GO-TO-MARKET ──
    slide = prs.slides.add_slide(blank)
    slide_header(slide, 6, "How We Win the Business",
                 "Concrete plan to land the first client and scale to a Wipro practice.")

    # 30-day plan
    box = add_shape_box(slide, 0.7, 2.2, 6.0, 3.5, border_color=WIPRO)
    tf = box.text_frame
    tf.word_wrap = True
    add_text(tf, "Land First Client (Next 30 Days)", size=14, bold=True, color=WIPRO_LT, space_after=10)

    steps_30 = [
        ("Week 1", "Polish prototype. Record 3-minute demo video of AI buying a product end-to-end."),
        ("Week 2", "Identify 5 target clients from existing Wipro accounts (SAP Commerce, D2C, digital-forward CxO)."),
        ("Week 3", "Send 1-pager + demo video. Frame as competitive urgency. Request 30-min meeting."),
        ("Week 4", "Live demo + this deck. Present timeline. Close first SOW."),
    ]
    for wk, desc in steps_30:
        add_text(tf, wk, size=11, bold=True, color=WHITE, space_after=1)
        add_text(tf, desc, size=10, color=GRAY, space_after=6)

    # What clinches the deal
    box = add_shape_box(slide, 7.1, 2.2, 5.5, 3.5, border_color=ACCENT)
    tf = box.text_frame
    tf.word_wrap = True
    add_text(tf, "What Clinches the Deal", size=14, bold=True, color=ACCENT, space_after=10)
    clinchers = [
        "Live demo — 'Look, Gemini is buying your product right now'",
        "Competitive fear — 'Samsung already supports this'",
        "Low risk — 'We don't touch your SAP. Minimal infra cost.'",
        "Fast delivery — '10 weeks to production'",
        "Working prototype — 'We're not selling PowerPoint'",
    ]
    for c in clinchers:
        add_bullet(tf, c, size=11, color=GRAY, space_after=5)

    # Scale plan
    phases = [
        ("Phase 2: Scale (6-12 Months)", WIPRO_LT, [
            "Sell to 5-10 more retail/CPG clients",
            "Build connector library (SAP, SFCC, Shopify, Magento)",
            "Case study + webinar from first client",
            "Cross-sell to every Wipro retail account",
        ]),
        ("Phase 3: Practice (12-24 Months)", WIPRO_LT, [
            "Dedicated team & P&L",
            "IP: accelerator kit, connectors, test suites",
            "Partnerships with Google, Shopify, Stripe",
            "NRF / Shoptalk presentations. Analyst briefings.",
        ]),
    ]
    for i, (title, clr, items) in enumerate(phases):
        x = 0.7 + i * 6.3
        box = add_shape_box(slide, x, 6.0, 5.95, 1.4, border_color=clr)
        tf = box.text_frame
        tf.word_wrap = True
        add_text(tf, title, size=12, bold=True, color=clr, space_after=4)
        for item in items:
            add_bullet(tf, item, size=10, color=GRAY, space_after=1)

    # ── SLIDE 7: COMPETITIVE ──
    slide = prs.slides.add_slide(blank)
    slide_header(slide, 7, "Competitive Landscape & Next Steps",
                 "No SI has a packaged agentic commerce offering. First-mover advantage is real.")

    comp_data = [
        ["Competitor", "Their Position", "Wipro's Advantage"],
        ["Accenture", "No known agentic commerce offering", "We have a working prototype. First-mover."],
        ["Infosys", "No known offering", "We have SAP Commerce expertise + demo ready"],
        ["TCS", "No known offering", "We move faster. 10 weeks vs their 6-month cycles."],
        ["Shopify (native)", "Building native UCP", "Only helps Shopify clients. SAP/SFCC clients need us."],
        ["SAP (native)", "Roadmap: 12-18 months out", "Clients need it NOW. We bridge the gap."],
    ]
    add_table(slide, 0.7, 2.2, 11.9, comp_data, col_widths=[2.5, 4.5, 4.9])

    # Next steps
    box = add_shape_box(slide, 0.7, 4.6, 11.9, 2.6, border_color=WIPRO)
    tf = box.text_frame
    tf.word_wrap = True
    add_text(tf, "Next Steps — What We Need from Leadership", size=16, bold=True, color=WIPRO_LT, space_after=12)
    next_steps = [
        "Approve investment: 2-3 developers for 10 weeks (existing resources, no new hires)",
        "Identify first client: Dyson, or similar SAP Commerce client in our portfolio",
        "Demo to client: Within 2 weeks (prototype is ready)",
        "First SOW signed: Target within 30 days",
        "First client live: 10 weeks after SOW",
    ]
    for ns in next_steps:
        add_bullet(tf, ns, size=13, color=GRAY, space_after=5)

    # CTA
    box = add_shape_box(slide, 2.5, 7.3, 8.3, 0.01)  # invisible spacer

    # ── SLIDE 8: Q&A ──
    slide = prs.slides.add_slide(blank)
    slide_header(slide, 8, "Anticipated Leadership Questions",
                 "Prepared answers to the questions that will come up.")

    qa1_data = [
        ["Market & Strategy", "Answer"],
        ["Is this a real trend or hype?", "Google, OpenAI, Microsoft all launched within weeks. Walmart, Visa, Mastercard backing it. This is real."],
        ["Will one protocol win?", "Too early. Our approach is protocol-agnostic — we support all. Like building a website that works on all browsers."],
        ["What if this doesn't take off?", "Low risk. 2-3 developers for 10 weeks. If it works, we have a practice. If not, skills transfer elsewhere."],
        ["Why would clients pay us vs DIY?", "Same reason they pay us for SAP. Protocol knowledge + pre-built templates save months. 10 weeks vs 6+ months internal."],
    ]
    add_table(slide, 0.7, 2.2, 11.9, qa1_data, col_widths=[3.5, 8.4])

    qa2_data = [
        ["Technical (Simplified)", "Answer"],
        ["Do we integrate with each AI separately?", "No. We expose endpoints, they call us. We register on portals (form filling). Code is 90% shared."],
        ["Is this secure?", "Payment NEVER touches our layer. Google Pay/Stripe handle it. No credit cards, no PII stored."],
        ["Do we change anything in SAP?", "Zero changes. We only call existing REST APIs. No config, no logic, no data changes."],
        ["Why not MuleSoft?", "If client has it, we use it. Otherwise $100K/yr license for 5 endpoints is overkill."],
    ]
    add_table(slide, 0.7, 4.5, 11.9, qa2_data, col_widths=[3.5, 8.4])

    return prs


# ── Main ──
if __name__ == "__main__":
    out_dir = "/Users/rohitgupta/Desktop/Claude/sap-ucp-poc"

    print("Creating Deck 1: Technical Leadership...")
    prs1 = create_deck1()
    path1 = os.path.join(out_dir, "Wipro-Agentic-Commerce-Technical.pptx")
    prs1.save(path1)
    print(f"  Saved: {path1}")

    print("Creating Deck 2: Business Leadership...")
    prs2 = create_deck2()
    path2 = os.path.join(out_dir, "Wipro-Agentic-Commerce-Leadership.pptx")
    prs2.save(path2)
    print(f"  Saved: {path2}")

    print("\nDone! Two PPTX files created.")
