#!/usr/bin/env python3
"""Generate a 5-slide demo deck for leadership presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette (Light Theme) ──
BG_DARK = RGBColor(0xFF, 0xFF, 0xFF)       # White background
BG_CARD = RGBColor(0xF5, 0xF5, 0xF8)       # Light gray card
WIPRO = RGBColor(0x5A, 0x2F, 0xC2)         # Wipro purple
WIPRO_LT = RGBColor(0x6B, 0x3F, 0xD4)      # Slightly lighter purple
ACCENT = RGBColor(0xD3, 0x00, 0x0F)        # Red (slightly darker for light bg)
GREEN = RGBColor(0x2E, 0x7D, 0x32)         # Darker green for readability
BLUE = RGBColor(0x1E, 0x88, 0xE5)          # Darker blue for readability
ORANGE = RGBColor(0xE6, 0x7E, 0x00)        # Darker orange
CYAN = RGBColor(0x00, 0x97, 0xA7)          # Darker cyan
WHITE = RGBColor(0x1A, 0x1A, 0x2E)         # Dark text (replaces white)
GRAY = RGBColor(0x4A, 0x4A, 0x5A)          # Dark gray text
MUTED = RGBColor(0x78, 0x78, 0x88)         # Muted text
GOLD = RGBColor(0xC4, 0x9A, 0x00)          # Darker gold for readability
BORDER = RGBColor(0xD8, 0xD8, 0xE0)        # Light border


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def tb(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))


def txt(tf, text, size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT, after=4, font="Calibri"):
    p = tf.add_paragraph() if tf.paragraphs[0].text != "" else tf.paragraphs[0]
    if tf.paragraphs[0].text != "" and p == tf.paragraphs[0]:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    p.space_after = Pt(after)
    return p


def bullet(tf, text, size=11, color=GRAY, after=3, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.font.bold = bold
    p.space_after = Pt(after)
    return p


def box(slide, left, top, width, height, fill=BG_CARD, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def mono(tf, text, size=8, color=GRAY, after=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = "Courier New"
    p.space_after = Pt(after)
    return p


def divider(slide, left, top, width, color=MUTED):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False


def add_table(slide, left, top, width, rows_data, col_widths=None, rh=0.3):
    rows, cols = len(rows_data), len(rows_data[0])
    ts = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(rh * rows))
    tbl = ts.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(9 if r == 0 else 10)
            p.font.bold = (r == 0)
            p.font.color.rgb = WIPRO_LT if r == 0 else GRAY
            p.font.name = "Calibri"
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xF0) if r == 0 else RGBColor(0xFF, 0xFF, 0xFF)
    return ts


def header(slide, title, subtitle=""):
    set_slide_bg(slide, BG_DARK)
    # Wipro badge
    s = box(slide, 12.1, 0.25, 0.5, 0.5, fill=WIPRO)
    s.text_frame.paragraphs[0].text = "W"
    s.text_frame.paragraphs[0].font.size = Pt(16)
    s.text_frame.paragraphs[0].font.bold = True
    s.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Title
    t = tb(slide, 0.6, 0.25, 10, 0.6)
    txt(t.text_frame, title, size=24, bold=True, color=WHITE, after=0)
    if subtitle:
        t = tb(slide, 0.6, 0.85, 11.5, 0.4)
        txt(t.text_frame, subtitle, size=11, color=MUTED, after=0)
    divider(slide, 0.6, 1.25, 12.1, WIPRO)


def create_demo_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ================================================================
    #  SLIDE 1: What Is UCP & What It Does
    # ================================================================
    slide = prs.slides.add_slide(blank)
    header(slide,
           "What Is Universal Commerce Protocol (UCP)?",
           "An open standard launched by Google (Jan 2026) that lets AI agents shop on behalf of users — browse, cart, checkout, track — all via conversation.")

    # LEFT: What Happened
    b = box(slide, 0.6, 1.45, 5.9, 2.0, border=ACCENT)
    tf = b.text_frame; tf.word_wrap = True
    txt(tf, "The AI Commerce Wave (Jan-Feb 2026)", size=13, bold=True, color=WHITE, after=6)
    for title, desc, clr in [
        ("Google — Universal Commerce Protocol (UCP)", "Buy inside Gemini & Google Search via Google Pay", ACCENT),
        ("OpenAI — Agent Commerce Protocol (ACP)", '"Buy it in ChatGPT" with Stripe — 500M+ users', BLUE),
        ("Microsoft — Copilot Checkout", "Shopping inside Windows/Edge/Office via Shopify", WIPRO_LT),
    ]:
        bullet(tf, title, size=10, color=clr, after=0, bold=True)
        bullet(tf, desc, size=9, color=GRAY, after=5)

    # RIGHT: How UCP Works
    b = box(slide, 6.8, 1.45, 5.9, 2.0, border=WIPRO)
    tf = b.text_frame; tf.word_wrap = True
    txt(tf, "How UCP Works", size=13, bold=True, color=WIPRO_LT, after=6)
    for s in [
        '1.  Merchant exposes a few REST endpoints (cart, checkout)',
        '2.  Publishes a discovery profile at /.well-known/ucp',
        '3.  AI platforms (Gemini, ChatGPT) discover & call these endpoints',
        '4.  User says "Buy me X" — AI handles the entire purchase flow',
        '5.  Payment via Google Pay / Stripe — never touches merchant',
    ]:
        bullet(tf, s, size=9, color=GRAY, after=2)

    # Era Table
    add_table(slide, 0.6, 3.7, 7.5,
              [["Era", "How Customers Shop", "What Merchants Need", "Standard"],
               ["2000s", "Type URLs in browser", "A website", "HTTP / HTML"],
               ["2010s", "Tap apps on phone", "A mobile app", "REST APIs"],
               ["2020s", "Search marketplaces", "Marketplace presence", "Amazon / Flipkart"],
               ["2026+", 'Tell AI: "Buy me a Dyson V15"', "AI-compatible commerce APIs", "UCP / ACP / MCP"]],
              col_widths=[1.0, 2.2, 2.3, 2.0], rh=0.3)

    # Who's Already In
    b = box(slide, 8.4, 3.7, 4.3, 1.5, border=GREEN)
    tf = b.text_frame; tf.word_wrap = True
    txt(tf, "Already Onboard (20+ Partners)", size=11, bold=True, color=GREEN, after=4)
    bullet(tf, "Retailers: Walmart, Target, Best Buy, Etsy, Wayfair, Macy's, Flipkart, Zalando", size=9, color=GRAY, after=2)
    bullet(tf, "Payments: Visa, Mastercard, Stripe, Adyen", size=9, color=GRAY, after=2)
    bullet(tf, "Platforms: Shopify, Google, OpenAI, Microsoft", size=9, color=GRAY, after=0)

    # Key Insight
    b = box(slide, 0.6, 5.5, 12.1, 1.2, border=ACCENT)
    tf = b.text_frame; tf.word_wrap = True
    txt(tf, "KEY INSIGHT", size=9, bold=True, color=ACCENT, after=3)
    txt(tf, "UCP is NOT an API that Google provides. Merchants build the API. Google's Gemini calls it.", size=13, bold=True, color=WHITE, after=3)
    txt(tf, 'Like building a website — you don\'t "integrate with Chrome." You build the site, browsers find you. Same here — expose endpoints, AI agents find you.', size=10, color=GRAY, after=0)

    # ================================================================
    #  SLIDE 2: Why This Is Needed
    # ================================================================
    slide = prs.slides.add_slide(blank)
    header(slide,
           "Why This Is Urgently Needed",
           "If a merchant isn't connected to AI shopping, they become invisible. AI will recommend competitors instead.")

    # WITHOUT
    b = box(slide, 0.6, 1.45, 5.9, 2.2, border=ACCENT)
    tf = b.text_frame; tf.word_wrap = True
    txt(tf, "WITHOUT Agentic Commerce", size=13, bold=True, color=ACCENT, after=6)
    for q, a in [
        ('Customer: "Hey Gemini, buy me a Dyson V15"', 'Dyson not connected. AI says: "I found similar from Samsung and Shark."'),
        ('Customer: "Reorder my Nespresso pods"', 'Nestle not connected. AI redirects to Amazon — lost margin, lost data.'),
        ('Customer: "Find me running shoes under $150"', 'Nike not connected. AI shows Adidas, New Balance, Puma instead.'),
    ]:
        txt(tf, q, size=9, bold=True, color=WHITE, after=0)
        txt(tf, "  " + a, size=9, color=RGBColor(0xC6, 0x28, 0x28), after=5)

    # WITH
    b = box(slide, 6.8, 1.45, 5.9, 2.2, border=GREEN)
    tf = b.text_frame; tf.word_wrap = True
    txt(tf, "WITH Agentic Commerce", size=13, bold=True, color=GREEN, after=6)
    for q, a in [
        ('Customer: "Hey Gemini, buy me a Dyson V15"', 'Dyson IS connected. AI shows V15, checkout in 30 sec via Google Pay. DIRECT SALE.'),
        ('Customer: "Reorder my Nespresso pods"', 'Nestle IS connected. AI knows their account, reorders same pods. RECURRING REVENUE.'),
        ('Customer: "Find me running shoes under $150"', 'Nike IS connected. AI shows Nike options first. FULL MARGIN, OWN DATA.'),
    ]:
        txt(tf, q, size=9, bold=True, color=WHITE, after=0)
        txt(tf, "  " + a, size=9, color=RGBColor(0x2E, 0x7D, 0x32), after=5)

    # Market Numbers
    for i, (num, label, clr) in enumerate([
        ("$385B", "Projected AI commerce\nmarket by 2028", GOLD),
        ("500M+", "ChatGPT users who\ncan now shop via AI", BLUE),
        ("2B+", "Google users with\nGemini AI shopping", ACCENT),
        ("20+", "Global partners already\nonboard with UCP", GREEN),
        ("90%", "Code reuse across\nUCP / ACP / MCP", WIPRO_LT),
    ]):
        x = 0.6 + i * 2.5
        b = box(slide, x, 3.9, 2.3, 1.1, border=clr)
        tf = b.text_frame; tf.word_wrap = True
        txt(tf, num, size=22, bold=True, color=clr, align=PP_ALIGN.CENTER, after=1)
        txt(tf, label, size=8, color=GRAY, align=PP_ALIGN.CENTER, after=0)

    # Threat
    b = box(slide, 0.6, 5.25, 12.1, 1.1, border=ORANGE)
    tf = b.text_frame; tf.word_wrap = True
    txt(tf, "THE COMPETITIVE THREAT", size=9, bold=True, color=ORANGE, after=3)
    txt(tf, "Walmart, Target, Best Buy, Shopify merchants are ALREADY connected. Every day a client waits, AI sends their customers to competitors.", size=12, bold=True, color=WHITE, after=3)
    txt(tf, "This is not a future trend — it launched in January 2026. The window to be a first-mover is NOW.", size=10, color=GRAY, after=0)

    # ================================================================
    #  SLIDE 3: Architecture Diagram (editable shapes)
    # ================================================================
    slide = prs.slides.add_slide(blank)
    header(slide,
           "Architecture — Where Wipro Sits",
           "Wipro builds a thin translation layer + optional platform extension. Everything else stays untouched.")

    # Helper: connector arrow (vertical)
    def varrow(slide, x, y, length, color=MUTED):
        shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(0.3), Inches(length))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.shadow.inherit = False

    # Helper: small label next to arrow
    def arrow_label(slide, x, y, text, color=MUTED):
        t = tb(slide, x, y, 2.5, 0.3)
        txt(t.text_frame, text, size=8, bold=True, color=color, after=0)

    # ── LAYER 1: END USER ──
    b = box(slide, 4.5, 1.45, 4.3, 0.55, fill=BG_CARD, border=BORDER)
    tf = b.text_frame; tf.word_wrap = True
    txt(tf, 'END USER  —  "Hey Gemini, buy me a Dyson V15"', size=10, bold=True, color=WHITE, after=0)

    varrow(slide, 6.5, 2.05, 0.35, BLUE)

    # ── LAYER 2: AI SHOPPING PLATFORMS ──
    b = box(slide, 0.6, 2.5, 9.5, 1.05, fill=BG_CARD, border=BLUE)
    tf = b.text_frame; tf.word_wrap = True
    txt(tf, "AI SHOPPING PLATFORMS  (not our scope)", size=11, bold=True, color=BLUE, after=2)
    txt(tf, "Google provides: user identity, Google Pay, shipping address, Shopping Graph, conversation context", size=8, color=MUTED, after=0)

    # Platform sub-boxes
    for i, (name, protocol, clr) in enumerate([
        ("Google Gemini", "UCP", BLUE),
        ("OpenAI ChatGPT", "ACP", RGBColor(0x6B, 0x3F, 0xD4)),
        ("Microsoft Copilot", "Checkout", RGBColor(0x00, 0x78, 0xD4)),
    ]):
        bx = box(slide, 0.8 + i * 3.2, 3.0, 2.9, 0.4, fill=BG_DARK, border=clr)
        tf = bx.text_frame; tf.word_wrap = True
        txt(tf, f"{name}  ({protocol})", size=8, bold=True, color=clr, after=0, align=PP_ALIGN.CENTER)

    varrow(slide, 6.5, 3.5, 0.35, WIPRO)
    arrow_label(slide, 6.9, 3.5, "HTTPS REST calls", WIPRO)

    # ── LAYER 3: WIPRO ACCELERATOR ──
    b = box(slide, 0.6, 3.95, 9.5, 1.5, fill=BG_CARD, border=WIPRO)
    tf = b.text_frame; tf.word_wrap = True
    txt(tf, "WIPRO AGENTIC COMMERCE ACCELERATOR  (deployed on client's cloud)", size=11, bold=True, color=WIPRO_LT, after=2)
    txt(tf, "Lambda / Container / Cloud Run  |  Stateless  |  Session-to-cart mapping in Redis", size=8, color=MUTED, after=0)

    # Component sub-boxes
    components = [
        ("Discovery", "/.well-known/ucp"),
        ("Checkout APIs", "POST/GET/PUT\ncomplete/cancel"),
        ("Field Translation", "UCP <> SAP/SFCC\nfield mapping"),
        ("Identity & Auth", "OAuth2 authorize\ntoken, Bearer pass"),
        ("Order Webhooks", "JWT signed push\nto Google"),
    ]
    for i, (title, desc) in enumerate(components):
        bx = box(slide, 0.8 + i * 1.9, 4.55, 1.75, 0.75, fill=BG_DARK, border=WIPRO)
        tf = bx.text_frame; tf.word_wrap = True
        txt(tf, title, size=8, bold=True, color=WIPRO_LT, after=1, align=PP_ALIGN.CENTER)
        txt(tf, desc, size=6, color=GRAY, after=0, align=PP_ALIGN.CENTER, font="Courier New")

    varrow(slide, 6.5, 5.5, 0.3, ORANGE)
    arrow_label(slide, 6.9, 5.5, "Calls platform APIs", ORANGE)

    # ── LAYER 4: WIPRO PLATFORM EXTENSION ──
    b = box(slide, 0.6, 5.9, 9.5, 0.8, fill=BG_CARD, border=ORANGE)
    tf = b.text_frame; tf.word_wrap = True
    txt(tf, "WIPRO PLATFORM EXTENSION  (only if platform lacks required APIs)", size=10, bold=True, color=ORANGE, after=0)

    # Extension sub-boxes
    extensions = [
        ("SAP Commerce", "Custom OCC Ext"),
        ("Salesforce SFCC", "Custom Cartridge"),
        ("Shopify", "Custom App"),
        ("Magento", "Custom Module"),
    ]
    for i, (platform, ext) in enumerate(extensions):
        bx = box(slide, 0.8 + i * 2.4, 6.3, 2.2, 0.35, fill=BG_DARK, border=ORANGE)
        tf = bx.text_frame; tf.word_wrap = True
        txt(tf, f"{platform}: {ext}", size=7, bold=True, color=ORANGE, after=0, align=PP_ALIGN.CENTER)

    varrow(slide, 6.5, 6.7, 0.25, GREEN)

    # ── LAYER 5: CLIENT PLATFORM ──
    b = box(slide, 0.6, 7.0, 9.5, 0.4, fill=BG_CARD, border=GREEN)
    tf = b.text_frame; tf.word_wrap = True
    txt(tf, "CLIENT'S COMMERCE PLATFORM (untouched)  —  Pricing, Tax, Inventory, Fulfillment, All Business Logic", size=9, bold=True, color=GREEN, after=0)

    # ── RIGHT SIDE: Payment Flow ──
    b = box(slide, 10.4, 2.5, 2.2, 1.6, fill=BG_CARD, border=GOLD)
    tf = b.text_frame; tf.word_wrap = True
    txt(tf, "PAYMENT", size=10, bold=True, color=GOLD, after=2)
    txt(tf, "(separate — never touches Wipro)", size=7, color=MUTED, after=4)
    for name in ["Google Pay", "Stripe / Adyen"]:
        bullet(tf, name, size=8, color=GOLD, after=2)
    txt(tf, "Tokenized, PCI-compliant", size=7, color=MUTED, after=0)

    # ── RIGHT SIDE: Order Updates ──
    b = box(slide, 10.4, 4.3, 2.2, 1.7, fill=BG_CARD, border=CYAN)
    tf = b.text_frame; tf.word_wrap = True
    txt(tf, "ORDER UPDATES", size=10, bold=True, color=CYAN, after=2)
    txt(tf, "(Webhook: Merchant > Google)", size=7, color=MUTED, after=4)
    for step in [
        "SAP fires status event",
        "Wipro builds full order",
        "Signs with JWT",
        "POST to Google webhook",
        "User gets update",
    ]:
        bullet(tf, step, size=7, color=GRAY, after=1)

    # ── RIGHT SIDE: What We Don't Build ──
    b = box(slide, 10.4, 6.2, 2.2, 1.2, fill=BG_CARD, border=ACCENT)
    tf = b.text_frame; tf.word_wrap = True
    txt(tf, "NOT OUR SCOPE", size=9, bold=True, color=ACCENT, after=3)
    for item in ["A2A Protocol", "ADK / AI Agent", "Chatbot / LLM", "Payment processing"]:
        bullet(tf, item, size=7, color=ACCENT, after=1)

    # ================================================================
    #  SLIDE 4: What Wipro Builds, Progress & Plan (UCP focused)
    # ================================================================
    slide = prs.slides.add_slide(blank)
    header(slide,
           "Wipro's Scope, Progress & Plan — UCP for Google",
           "What we build, what we don't, what's done, what's remaining, and the step-by-step plan.")

    # ── TOP LEFT: What Wipro Builds ──
    b = box(slide, 0.6, 1.45, 4.0, 2.7, border=WIPRO)
    tf = b.text_frame; tf.word_wrap = True
    txt(tf, "What Wipro Builds", size=12, bold=True, color=WIPRO_LT, after=4)
    for title, desc in [
        ("UCP Discovery Profile", "/.well-known/ucp — declares capabilities, payment, signing keys"),
        ("5 Checkout REST Endpoints", "POST/GET/PUT checkout-sessions, complete, cancel"),
        ("Field Translation Layer", "UCP fields <> SAP OCC fields (city > town, cents > dollars)"),
        ("OAuth2 Identity Linking", "/authorize + /token — links Google account to merchant"),
        ("Order Status Webhooks", "Full order push to Google, signed with JWT (UCP profile keys)"),
        ("Platform Extension", "SAP OCC addon / SFCC cartridge — only if platform lacks APIs"),
    ]:
        txt(tf, title, size=9, bold=True, color=WIPRO_LT, after=0)
        bullet(tf, desc, size=7, color=GRAY, after=3)

    # ── TOP MIDDLE: What Wipro Does NOT Build ──
    b = box(slide, 4.8, 1.45, 3.6, 2.7, border=ACCENT)
    tf = b.text_frame; tf.word_wrap = True
    txt(tf, "What Wipro Does NOT Build", size=12, bold=True, color=ACCENT, after=4)
    for title, desc in [
        ("A2A Protocol", "Google's internal agent-to-agent comms — their plumbing"),
        ("ADK / AI Agent", "Google builds Gemini AI — we just receive its REST calls"),
        ("Chatbot / LLM", "All AI reasoning & conversation is Google's scope"),
        ("Payment Processing", "Google Pay token > Stripe/Adyen decrypts & charges"),
        ("Business Logic", "Pricing, tax, inventory, promotions — stays in SAP"),
    ]:
        txt(tf, "X  " + title, size=9, bold=True, color=ACCENT, after=0)
        bullet(tf, desc, size=7, color=GRAY, after=3)

    # ── TOP RIGHT: What We've Done So Far ──
    b = box(slide, 8.6, 1.45, 4.1, 2.7, border=GREEN)
    tf = b.text_frame; tf.word_wrap = True
    txt(tf, "What We've Done So Far", size=12, bold=True, color=GREEN, after=4)
    for title, desc in [
        ("Working POC", "Spring Boot app — all 5 UCP endpoints + discovery + SAP OCC mock"),
        ("Full UCP Spec Research", "Read complete Google UCP spec from ucp.dev — checkout, orders, identity"),
        ("Field Mapping (UCP > SAP)", "Mapped all fields: address, cart, checkout, order status"),
        ("Identity Linking", "OAuth2 authorize/token endpoints built and working in POC"),
        ("Architecture Design", "4-layer design: Google > Wipro UCP > Extension > Client Platform"),
    ]:
        txt(tf, title, size=9, bold=True, color=GREEN, after=0)
        bullet(tf, desc, size=7, color=GRAY, after=3)

    # ── BOTTOM LEFT: What's Still Needed ──
    b = box(slide, 0.6, 4.4, 5.5, 2.8, border=ORANGE)
    tf = b.text_frame; tf.word_wrap = True
    txt(tf, "What's Still Needed (UCP Only)", size=12, bold=True, color=ORANGE, after=4)
    for title, desc in [
        ("Connect to Real SAP", "Replace mock with actual SAP Commerce OCC APIs — need sandbox access"),
        ("Google Merchant Center", "Register merchant, upload product feed with native_commerce=TRUE"),
        ("Google Pay via PSP", "Set up Stripe/Adyen to accept Google Pay tokenized payment tokens"),
        ("Order Webhooks", "SAP status events > Wipro translates > signed JWT push to Google"),
        ("Security & Signing", "JWK key pair for request signing, OAuth2 hardening, HTTPS"),
        ("Testing & Monitoring", "End-to-end flow testing, latency checks, logging, error handling"),
    ]:
        txt(tf, title, size=9, bold=True, color=ORANGE, after=0)
        bullet(tf, desc, size=7, color=GRAY, after=2)

    # ── BOTTOM RIGHT: The Plan ──
    b = box(slide, 6.3, 4.4, 6.4, 2.8, border=BLUE)
    tf = b.text_frame; tf.word_wrap = True
    txt(tf, "The Plan — UCP (Google/Gemini)", size=12, bold=True, color=BLUE, after=4)

    for step, what, duration, clr in [
        ("Step 1: Client + SAP Sandbox", "Identify SAP Commerce client. Get OCC sandbox access.", "Week 1", WIPRO_LT),
        ("Step 2: Connect to Real SAP", "Replace mock with real OCC APIs. Validate field mappings.", "Week 2-3", BLUE),
        ("Step 3: Google Merchant + Pay", "Register in Merchant Center, upload feed, set up Google Pay PSP.", "Week 3-4", GREEN),
        ("Step 4: Order Webhooks + Signing", "SAP events > Wipro > signed JWT push to Google webhook_url.", "Week 4-5", ORANGE),
        ("Step 5: Identity Linking", "OAuth2 flow with real SAP customer accounts. Guest + linked.", "Week 5-6", CYAN),
        ("Step 6: E2E Test + Go-Live", "Full flow: Gemini > Wipro > SAP > order confirmed. Security review.", "Week 7-8", GOLD),
    ]:
        txt(tf, f"{step}  ({duration})", size=9, bold=True, color=clr, after=0)
        bullet(tf, what, size=7, color=GRAY, after=2)

    return prs


if __name__ == "__main__":
    out_dir = "/Users/rohitgupta/Desktop/Claude/sap-ucp-poc"
    print("Creating 4-Slide Demo Deck...")
    prs = create_demo_deck()
    path = os.path.join(out_dir, "Wipro-Agentic-Commerce-Demo.pptx")
    prs.save(path)
    print(f"Saved: {path}")
    print("Done!")
